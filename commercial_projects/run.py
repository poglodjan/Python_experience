from pptx import Presentation
import tkinter as tk
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from datetime import datetime
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from charts import *
from scenarios import *
from datetime import datetime
import os
from io import BytesIO

############################################
#
# 1) Creating first slides
#
############################################

def run_pp(settings_frame,selected_options,selected_inputs, is_importer, entry_boxes):
    counter=0
    loading_label = tk.Label(settings_frame, text="Loading: 0%",background="#404040",
                             font=("Arial",16),foreground="#8BE753")
    loading_label.grid(row=3, column=4, columnspan=2, padx=5, sticky=tk.E)

    try: presentation = Presentation("template.pptx")
    except Exception as e: print("error")
    is_importer=is_importer.get()
    print_options=""
    for i in range(len(selected_options)):
        print_options+= f"0{i} " + selected_options[i] + '\n'
    base = entry_boxes[0][0].get()
    term = entry_boxes[0][1].get()
    notional = float(''.join(entry_boxes[0][2].get().split()))
    expiry_date = entry_boxes[0][3].get()
    spot = float(entry_boxes[1][0].get()) 
    company = entry_boxes[1][1].get()
    s_min = float(entry_boxes[2][0].get())
    s_max = float(entry_boxes[2][1].get())
    s_decimal = int(entry_boxes[2][2].get())
    c_min = float(entry_boxes[3][0].get())
    c_max = float(entry_boxes[3][1].get())
    c_decimal = int(entry_boxes[3][2].get())
    if is_importer: header_name = ["Spot at expiry", "Buying rate in structure", f"Outcome on transaction in {term}"]
    if not is_importer: header_name = ["Spot at expiry", "Selling rate in structure", f"Outcome on transaction in {term}"]

    # setting textes on slide _1_
    slide = presentation.slides[0]
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(8)), top=Inches(to_cm(0.5)), width=Inches(to_cm(20)), height=Inches(to_cm(1.18)))
    text_frame = textbox.text_frame
    text_frame.text = f"{datetime.today().strftime('%d-%m-%Y')} | Market Commentary Intended for Institutional Clients Only"
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.name = "Citi Sans Text"
    
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(1.84)), top=Inches(to_cm(10.2)), width=Inches(to_cm(25.4)), height=Inches(to_cm(2.3)))
    text_frame = textbox.text_frame
    text_frame.text = f"Presentation for {company}"
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(50)
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.name = "Citi Sans Display"
    
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(1.84)), top=Inches(to_cm(13.8)), width=Inches(to_cm(25.4)), height=Inches(to_cm(2.3)))
    text_frame = textbox.text_frame
    text_frame.text = f"Hedging {base}/{term} exposure"
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(20)
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.name = "Citi Sans Display"
    
     # setting textes on slide _4_
    slide = presentation.slides[3]
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(9.8)), top=Inches(to_cm(7.33)), width=Inches(to_cm(14.6)), height=Inches(to_cm(1.8)))
    text_frame = textbox.text_frame
    text_frame.text = f"{base}/{term} Risk - Executive Summary"
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(24)
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.name = "Citi Sans Display"
    
     # setting textes on slide _5_
    slide = presentation.slides[4]
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(5.06)), top=Inches(to_cm(8.54)), width=Inches(to_cm(8.44)), height=Inches(to_cm(1.36)))
    text_frame = textbox.text_frame
    text_frame.text = f"{'{:,.0f}'.format(notional)} {base}"
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(20)
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.name = "Citi Sans Text"

    textbox = slide.shapes.add_textbox(left=Inches(to_cm(5.86)), top=Inches(to_cm(12.84)), width=Inches(to_cm(8.17)), height=Inches(to_cm(1.36)))
    text_frame = textbox.text_frame
    text_frame.text = f"'{float(spot):.{s_decimal}f}'"
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(20)
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.name = "Citi Sans Text"
    
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(5.44)), top=Inches(to_cm(10.64)), width=Inches(to_cm(7.74)), height=Inches(to_cm(1.57)))
    text_frame = textbox.text_frame
    text_frame.text = f"{expiry_date}"
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(20)
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.name = "Citi Sans Text"
    
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(14.55)), top=Inches(to_cm(8.01)), width=Inches(to_cm(11.85)), height=Inches(to_cm(8.19)))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = f"{print_options}"
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255,255,255)
        paragraph.font.name = "Citi Sans Text" 
    
    if is_importer:
        textbox = slide.shapes.add_textbox(left=Inches(to_cm(0.63)), top=Inches(to_cm(2.76)), width=Inches(to_cm(32.25)), height=Inches(to_cm(1.8)))
        text_frame = textbox.text_frame
        text_frame.text = f"Citi is pleased to present {base}/{term} risk management solutions for {company}\n {company} is buying {base} against {term}"
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = RGBColor(0, 32, 96)
            paragraph.font.name = "Citi Sans Text"
    if not is_importer:
        textbox = slide.shapes.add_textbox(left=Inches(to_cm(0.63)), top=Inches(to_cm(2.76)), width=Inches(to_cm(32.25)), height=Inches(to_cm(1.8)))
        text_frame = textbox.text_frame
        text_frame.text = f"Citi is pleased to present {base}/{term} risk management solutions for {company}\n {company} is selling {base} against {term}"
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = RGBColor(0, 32, 96)
            paragraph.font.name = "Citi Sans Text"
    

############################################
#
# 2) Creating slides of models
#
############################################

    # setting models and putting on slides
    if "FX Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    #data
        s = "FX Forward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        zp = 1/float(row[1].get())*100
        lever = float(row[1].get())
        array_ = forward_fun(strike,zp,is_importer,c_min,c_max,c_decimal)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 5 
        else: slide_n = 33

    # payout
        if lever!=1:
            if is_importer: tekst = f"• At settlement date, if spot is below strike, {company} will buy {base} Leveraged Notional at {float(strike):.{s_decimal}f}\n" \
            f"• If spot is above strike, {company} can buy {base} Notional at {float(strike):.{s_decimal}f}"
            if not is_importer: tekst = f"• At settlement date, if spot is below strike, {company} will sell {base} Notional at {float(strike):.{s_decimal}f}\n" \
            f"• If spot is above strike, {company} can sell {base} Leveraged Notional at {float(strike):.{s_decimal}f}"
        if lever==1:
            if is_importer: tekst = f"• At settlement date, {company} will buy {base} Notional at {float(strike):.{s_decimal}f}"
            if not is_importer: tekst = f"• At settlement date, {company} will sell {base} Notional at {float(strike):.{s_decimal}f}"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms =f"• {base}/{term} Spot: {float(spot):.{s_decimal}f}\n" + \
       f"• Expiry Date: {expiry_date}\n" + \
       f"• Strike: {float(strike):.{s_decimal}f}\n" + \
       f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if lever!=1: terms+=f"\n• Leveraged Notional: {'{:,.0f}'.format(notional)}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Structure without upfront cost providing known {base} buying rate",
        "• No participation in FX rate drop\n• Opportunity cost if spot is below Strike at expiry"]
        if not is_importer: pros_cons = [f"• Structure without upfront cost providing known {base} selling rate",
        "• No participation in FX rate increase\n• Opportunity cost if spot is above Strike at expiry"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Simplest FX hedging solution providing full hedge, but no participation in {base} rate weakening"
        if not is_importer: subtitle = f"Simplest FX hedging solution providing full hedge, but no participation in {base} rate streghtening"
        do_subtitle(presentation,subtitle,slide_n)
    
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close() 
        
    # scenario
        scenario = forward_scenario(strike, zp, notional,is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Parforward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Parforward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        zp = 1/float(row[1].get())*100
        array_ = forward_fun(strike,zp,is_importer,c_min,c_max,c_decimal)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 6 
        else: slide_n = 34
    
    # payout
        if lever!=1:
            if is_importer: tekst = f"• At settlement date, if spot is below strike, {company} will buy {base} Leveraged Notional at {float(strike):.{s_decimal}f}\n" \
            f"• If spot is above strike, {company} can buy {base} Notional at {float(strike):.{s_decimal}f}"
            if not is_importer: tekst = f"• At settlement date, if spot is below strike, {company} will sell {base} Notional at {float(strike):.{s_decimal}f}\n" \
            f"• If spot is above strike, {company} can sell {base} Leveraged Notional at {float(strike):.{s_decimal}f}"
        if lever==1:
            if is_importer: tekst = f"• At settlement date, {company} will buy {base} Notional at {float(strike):.{s_decimal}f}"
            if not is_importer: tekst = f"• At settlement date, {company} will sell {base} Notional at {float(strike):.{s_decimal}f}"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms =f"• {base}/{term} Spot: {float(spot):.{s_decimal}f}\n" + \
       f"• Expiry Date: {expiry_date}\n" + \
       f"• Strike: {float(strike):.{s_decimal}f}\n" + \
       f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if lever!=1: terms+=f"\n• Leveraged Notional: {'{:,.0f}'.format(notional)}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Structure without upfront cost providing known {base} buying rate",
        "• No participation in FX rate drop\n• Opportunity cost if spot is below Strike at expiry"]
        if not is_importer: pros_cons = [f"• Structure without upfront cost providing known {base} selling rate",
        "• No participation in FX rate increase\n• Opportunity cost if spot is above Strike at expiry"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"A strip of synthetic Forwards with single guaranteed {base} buying rate"
        if not is_importer: subtitle = f"A strip of synthetic Forwards with single guaranteed {base} selling rate"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    # scenario
        scenario = forward_scenario(strike, zp, notional,is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)
        
####################################################################################################################################

    if "European Call Option" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "European Call Option"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        koszt = float(row[1].get())
        nominal = float(''.join(row[2].get().split()))
        if not is_importer: s = "European Put Option"
        array_ = Eur_fun(strike, 100, koszt, spot, is_importer, c_min, c_max)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 7
        else: slide_n = 35
    
    # payout
        if is_importer: tekst = f"At each expiry date:\n" \
             f"• If {base}/{term} is above {strike:,.{s_decimal}f} {company} will buy {base} Notional at {strike:,.{s_decimal}f}\n" \
             f"• If {base}/{term} is at or below {strike:,.{s_decimal}f} {company} can buy {base} Notional at market rate"
        if not is_importer: tekst = f"At each expiry date:" \
             f"• If {base}/{term} is below {strike:,.{s_decimal}f} {company} will sell {base} Notional at {strike:,.{s_decimal}f} " \
             f"• If {base}/{term} is above {strike:,.{s_decimal}f} {company} can sell {base} Notional at market rate"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = f"• To fully hedge against unfavorable FX movements but retain full participation in favorable spot moves" \
             f"• Very often, in high yielding currencies, cost of buying {base} put option is cheaper than cost of carry in Forward" \
             f"• To take advantage of lower implied volatility, when it is cheap to buy the optionality " \
             f"• Accounting friendly"
        if not is_importer: whyto = f"• To fully hedge against unfavorable FX movements but retain full participation in favorable spot moves" \
             f"• Very often, in high yielding currencies, cost of selling {base} call option is cheaper than cost of carry in Forward" \
             f"• To take advantage of lower implied volatility, when it is cheap to buy the optionality " \
             f"• Can be used in more dynamic approach (e.g. if spot goes lower after the trade {company} can later sell {base} put to recover the premium paid upfront or restructure call into Forward rate off lower spot rate)" \
             f"• Accounting friendly"
    # terms
        terms = f"• {base}/{term} Spot: {spot:,.{s_decimal}f}\n" \
        f"• Expiry Date: {expiry_date}\n" \
        f"• Strike: {strike:,.{s_decimal}f}\n" \
        f"• Premium in %: {koszt:.{s_decimal}f}% {'{:,.0f}'.format(int(notional/lever))} {base}\n" \
        f"• Notional: {'{:,.0f}'.format(int(notional))}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Unlike Forward contract, option offers full participation in lower {base} buying rate\n• Full hedge above Strike level",
        f"• Upfront premium payment"]
        if not is_importer: pros_cons = [f"• Unlike Forward contract, option offers full participation in higher {base} selling rate\n• Full hedge below Strike level",
        f"• Upfront premium payment"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in lower spot move"
        if not is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in higher spot move"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    # scenario
        scenario = Eur_scenario(strike, nominal, is_importer, s_decimal,zp, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Call Spread" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Call Spread"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        koszt = float(row[2].get())
        nominal = float(''.join(row[3].get().split()))
        if not is_importer: s = "Put Spread"
        array_ = Spread_fun(lower,upper,koszt,spot,is_importer,c_min,c_max)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 8
        else: slide_n = 36
    
    # payout
        if is_importer: tekst = "At each expiry date:\n" + \
    f"• If spot is at or below {lower:,.{s_decimal}f} {company} can buy {base} Notional at market rate\n" + \
    f"• If spot is between {lower:,.{s_decimal}f} and {upper:,.{s_decimal}f} {company} will buy {base} Notional at {lower:,.{s_decimal}f}\n" + \
    f"• If spot is above {lower:,.{s_decimal}f}, {company} can effectively buy {base} Notional at market rate MINUS {upper - lower:.{s_decimal}f}" + \
    f"({upper:,.{s_decimal}f} - {lower:,.{s_decimal}f} = {upper - lower:,.{s_decimal}f})"
        if not is_importer: tekst = "At each expiry date:\n" + \
    f"• If spot is below {lower:,.{s_decimal}f} {company} will sell {base} Notional at market rate PLUS {upper - lower:,.{s_decimal}f}" + \
    f" ({upper:,.{s_decimal}f} - {lower:,.{s_decimal}f} = {upper - lower:,.{s_decimal}f})\n" + \
    f"• If spot is between {lower:,.{s_decimal}f} and {upper:,.{s_decimal}f} {company} will sell {base} Notional at {upper:,.{s_decimal}f}\n" + \
    f"• If spot is at or above {upper:,.{s_decimal}f} {company} can sell {base} Notional at market rate"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = f"• {base}/{term} Spot: {spot:,.{s_decimal}f} \n" \
        f"• Expiry Date: {expiry_date} \n" \
        f"• Lower Strike: {lower:,.{s_decimal}f}\n" \
        f"• Upper Strike: {upper:,.{s_decimal}f}\n" \
        f"• Premium in %: {koszt:,.{s_decimal}f}% {'{:,.0f}'.format(int(notional*koszt/100))} {base}\n" \
        f"• Notional: {'{:,.0f}'.format(int(notional))}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Full hedge between Lower and Upper Strikes\n• Unlike Forward contract, Call Spread offers full participation in {base} weakening\n• Opportunity to receive a subsidy equal to difference between Upper and Lower Strikes if spot is above Upper Strike",
        f"• Upfront premium payment\n• No full hedge above Upper Strike"]
        if not is_importer: pros_cons = [f"• Full hedge between Lower and Upper Strikes\n• Unlike Forward contract, Put Spread offers full participation in {base} strengthening\n• Opportunity to receive a subsidy equal to difference between Upper and Lower Strikes if spot is below Lower Strike",
        f"• Upfront premium payment\n• No full hedge below Lower Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Reduces the upfront premium by limiting protection range up to Upper Strike, while retaining full participation in lower spot moves"
        if not is_importer: subtitle = f"Reduces the upfront premium by limiting protection range down to Lower Strike, while retaining full participation in higher spot moves"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    #scenario
        scenario = Spread_scenario(lower,upper,koszt,nominal,is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Participating Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Participating Forward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        zp = float(row[1].get())
        array_ = Partyp_fun(strike,zp,is_importer,c_min,c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 9
        else: slide_n = 37
    
    # payout
        if is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is below {strike:.{s_decimal}f} {company}, " \
    f"will buy {zp}% of {base} Notional " \
    f"({(notional * (zp / 100)):.{s_decimal}f} {base} at {strike:.{s_decimal}f}, " \
    f"and remaining {100 - zp}% of {base} Notional " \
    f"({(notional - notional * (zp / 100)):.{s_decimal}f} {base} {company} can buy at market spot rate\n" \
    f"• If spot is at or above {strike:.{s_decimal}f}, " \
    f"{company} can buy 100% Notional at {strike:.{s_decimal}f}" 
        if not is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is at or above {strike:.{s_decimal}f} {company}, " \
    f"{company} will sell {zp}% of {base} Notional " \
    f"({(notional * (zp / 100)):.{s_decimal}f} {base} at {strike:.{s_decimal}f} and remaining {100 - zp}% of {base} Notional " \
    f"({(notional - notional * (zp / 100)):.{s_decimal}f} {base} {company} can sell at market spot rate\n" \
    f"• If spot is below {strike:.{s_decimal}f} {company}, " \
    f"{company} can sell 100% Notional at {strike:.{s_decimal}f}"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = f"• {base}/{term} Spot: {spot:.{s_decimal}f}\n" \
    f"• Expiry Date: {expiry_date}\n" \
    f"• Strike: {strike:.{s_decimal}f}\n" \
    f"• Participation: {(100 - zp)}%" \
    f"• Notional: {'{:,.0f}'.format(int(notional))}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Full hedge above Strike level\n• Partial participation in a favorable spot move",
        f"• Higher maximum {base} buying rate compared to Par Forward\n• Opportunity cost if spot at expiry is below Strike"]
        if not is_importer: pros_cons = [f"• Full hedge below Strike level\n• Partial participation in a favorable spot move",
        f"• Lower minimum {base} selling rate compared to Par Forward\n• Opportunity cost if spot at expiry is above Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Structure gives full protection with participation on {100 - zp}% of the notional in {term} strengthening"
        if not is_importer: subtitle = f"Structure gives full protection with participation on {100 - zp}% of the notional in {term} weakening"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    # scenario
        scenario = Partyp_scenario(strike, zp, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name, partyp=True)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Collar" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Collar"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        zp = 1/float(row[2].get())*100
        lever = float(row[2].get())
        array_ = kor_fun(lower,upper,zp,is_importer,c_min,c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 10
        else: slide_n = 38
    
        if lever==1:tek=""
        else: tek = "Leveraged"
    # payout
        if is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is at or below {lower:.{s_decimal}f} {company}, " \
    f"{company} will buy {base} {tek} Notional at {lower:.{s_decimal}f}\n" \
    f"• If spot is between {lower:.{s_decimal}f} and {upper:.{s_decimal}f} {company}, " \
    f"{company} can buy {base} Notional at prevailing market spot\n" \
    f"• If spot is at or above {upper:.{s_decimal}f} {company}, " \
    f"{company} will buy {base} Notional at {upper:.{s_decimal}f}"
        if not is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is at or below {lower:.{s_decimal}f} {company}, " \
    f"{company} will sell {base} Notional at {lower:.{s_decimal}f}\n" \
    f"• If spot is between {lower:.{s_decimal}f} and {upper:.{s_decimal}f} {company}, " \
    f"{company} can sell {base} Notional at prevailing market spot\n" \
    f"• If spot is at or above {upper:.{s_decimal}f} {company}, " \
    f"{company} will sell {base}tekst5 Notional at {upper:.{s_decimal}f}"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = f"• {base}/{term} Spot: {spot:.{s_decimal}f}\n" \
    f"• Expiry Date: {expiry_date}\n" \
    f"• Lower Strike: {lower:.{s_decimal}f}\n" \
    f"• Upper Strike: {upper:.{s_decimal}f}" \
    f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if lever!=1: terms+=f"\n• Leveraged Notional: {'{:,.0f}'.format(notional)}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Client knows the maximum {base} buying rate and can participate in rate decrease down to Lower Strike",
        f"• If spot goes below Lower Strike Client will be obliged to buy {base} at Lower Strike which will be higher than the prevailing spot market\n• Upper Strike is above Forward Rate"]
        if not is_importer: pros_cons = [f"• Client knows the minimum {base} selling rate and can participate in rate increase up to Upper Strike",
        f"• If spot goes above Upper Strike Client will be obliged to sell {base} at Upper Strike which will be lower than the prevailing spot market\n• Lower Strike is below Forward Rate"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Structure with no upfront premium cost offering full hedge above Upper Strike and at the same time allowing to participate in favorable spot move down to preagreed level"
        if not is_importer: subtitle = f"Structure with no upfront premium cost offering full hedge below Lower Strike and at the same time allowing to participate in favorable spot move up to preagreed level"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
        
    # scenario
        if lever ==1: scenario = Kor_scenario(lower, upper, zp, notional, is_importer, s_decimal, s_min, s_max)
        if lever !=1: scenario = Kor_scenario_lever(lower, upper, zp,lever, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################
        
    if "Capped Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Capped Forward"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        zp = 1/float(row[2].get())*100
        lever = float(row[2].get())
        if is_importer: array_ = kor_mewa_fun(lower, upper, lower, zp,
                           is_importer, c_min, c_max)
        if not is_importer: array_ = kor_mewa_fun(lower, upper, upper, zp,
                           is_importer, c_min, c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 11
        else: slide_n = 39

        if lever==1:tek=""
        else: tek = "Leveraged"
    # payout
        if is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is below {lower:.{s_decimal}f} {company}," \
    f"{company} will buy {base} {tek} Notional at {lower:.{s_decimal}f}\n" \
    f"• If spot is between {lower:.{s_decimal}f} and {upper:.{s_decimal}f} {company}," \
    f"{company} will buy {base} Notional at {lower:.{s_decimal}f}\n" \
    f"• If spot is at or above {upper:.{s_decimal}f}," \
    f"{company} can effectively buy {base} Notional at prevailing spot rate MINUS {round(upper - lower, 4)} {term}"
        if not is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is above {upper:.{s_decimal}f} {company}," \
    f"{company} will sell {base} {tek} Notional at {upper:.{s_decimal}f}\n" \
    f"• If spot is between {lower:.{s_decimal}f} and {upper:.{s_decimal}f} {company}, " \
    f"{company} will sell {base} Notional at {upper:.{s_decimal}f}\n" \
    f"• If spot is at or below {lower:.{s_decimal}f} {company}, " \
    f"{company} can effectively sell {base} Notional at prevailing spot rate PLUS {round(upper - lower, 4)} {term}"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        if is_importer:
            terms = f"• {base}/{term} Spot: {spot:.{s_decimal}f}\n" \
        f"• Expiry Date: {expiry_date}\n" \
        f"• Strike: {lower:.{s_decimal}f}\n" \
        f"• Upper Strike: {upper:.{s_decimal}f} \n" \
        f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        else:
            terms = f"• {base}/{term} Spot: {spot:.{s_decimal}f}\n" \
        f"• Expiry Date: {expiry_date}\n" \
        f"• Lower Strike: {lower:.{s_decimal}f}\n" \
        f"• Strike: {upper:.{s_decimal}f} \n" \
        f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if lever!=1: terms+=f"\n• Leveraged Notional: {'{:,.0f}'.format(notional)}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Initial {base} buying rate is lower compared to Forward\n• Lower credit charges than in Forward",
        f"• No hedge above Upper Strike\n• Opportunity cost if spot at expiry is below Strike"]
        if not is_importer: pros_cons = [f"• Initial {base} selling rate is higher compared to Forward\n• Lower credit charges than in Forward",
        f"• No hedge below Lower Strike\n• Opportunity cost if spot at expiry is above Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Allows to buy {base} at rate lower than Par Forward, but provides only partial hedge above Upper Strike"
        if not is_importer: subtitle = f"Allows to sell {base} at rate higher than Par Forward, but provides only partial hedge below Lower Strike"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()

    # scenario
        if zp == 100:
            if is_importer: scenario = kor_mewa_scenario(lower, upper, lower, zp, notional, is_importer, s_decimal, s_min, s_max)
            if not is_importer: scenario = kor_mewa_scenario(lower, upper, upper, zp, notional, is_importer, s_decimal, s_min, s_max)
        if zp!=100:
            if is_importer: scenario = kor_mewa_scenario_lever(lower, upper, lower, zp, lever, notional, is_importer, s_decimal, s_min, s_max)
            if not is_importer: scenario = kor_mewa_scenario_lever(lower, upper, upper, zp, lever, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Loss Capped Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Loss Capped Forward"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        zp = 1/float(row[2].get())*100
        lever = float(row[2].get())
        if is_importer: array_ = Kor_mewa_fun2(lower, upper, upper, zp,
                           is_importer, c_min, c_max)
        if not is_importer: array_ = Kor_mewa_fun2(lower, upper, lower, zp,
                           is_importer, c_min, c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 12
        else: slide_n = 40

        if lever==1:tek=""
        else: tek = "Leveraged"
    # payout
        if is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is below {lower:.{s_decimal}f} {company}, " \
    f"{company} will buy {base} {tek} Notional at market rate PLUS {round(upper - lower, 4)} {term}\n" \
    f"• If spot is between {lower:.{s_decimal}f} and {upper:.{s_decimal}f} {company}, " \
    f"{company} will buy {base} {tek} Notional at {upper:.{s_decimal}f}\n" \
    f"• If spot is at or above {upper:.{s_decimal}f} {company}, " \
    f"{company} will buy {base} Notional at {upper:.{s_decimal}f}"
        if not is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is above {upper:.{s_decimal}f} {company}, " \
    f"{company} will sell {base} {tek} Notional at market rate MINUS {round(upper - lower, 4)} {term}\n" \
    f"• If spot is between {lower:.{s_decimal}f} and {upper:.{s_decimal}f} {company}, " \
    f"{company} will sell {base} {tek} Notional at {lower:.{s_decimal}f}\n" \
    f"• If spot is at or below {lower:.{s_decimal}f} {company}, " \
    f"{company} will sell {base} Notional at {lower:.{s_decimal}f}"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        if is_importer:
            terms = f"• {base}/{term} Spot: {spot:.{s_decimal}f}\n" \
        f"• Expiry Date: {expiry_date}\n" \
        f"• Lower Strike: {lower:.{s_decimal}f}\n" \
        f"• Strike: {upper:.{s_decimal}f} \n" \
        f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        else:
            terms = f"• {base}/{term} Spot: {spot:.{s_decimal}f}\n" \
        f"• Expiry Date: {expiry_date}\n" \
        f"• Strike: {lower:.{s_decimal}f}\n" \
        f"• Upper Strike: {upper:.{s_decimal}f} \n" \
        f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if lever!=1: terms+=f"\n• Leveraged Notional: {'{:,.0f}'.format(notional)}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Provides full hedge\n• Maximum loss on the transaction is capped",
        f"• Higher {base} buying rate compared to Forward"]
        if not is_importer: pros_cons = [f"• Provides full hedge\n• Maximum loss on the transaction is capped",
        f"• Lower {base} selling rate compared to Forward"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Structure gives full protection with limited loss in case of {base} weakening"
        if not is_importer: subtitle = f"Structure gives full protection with limited loss in case of {base} strengthening"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    # scenario
        if is_importer: scenario = kor_mewa_scenario2(lower, upper, upper, zp, notional, is_importer, s_decimal, s_min, s_max)
        if not is_importer: scenario = kor_mewa_scenario2(lower, upper, lower, zp, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Step Down Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Step Down Forward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        wyp = -float(row[1].get())
        trigger = float(row[2].get())
        zp = 1/float(row[3].get())*100
        lever = float(row[3].get())
        if not is_importer: s = "Step Up Forward"
        array_ = forw_zrek_fun(strike, wyp, zp, trigger, is_importer, c_min, c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 13
        else: slide_n = 41

        if lever==1:tek=""
        else: tek = "Leveraged"
    # payout
        # scenario 1 ___
        if is_importer and trigger <= strike:
            if zp == 100:
                tekst = (
                    f"At each expiry date:\n"
                    f"• If spot is at or below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} Notional at {((strike - wyp / 100)):.{s_decimal}f}\n"
                    f"• If spot is above {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} Notional at {strike:.{s_decimal}f}"
                )
            elif zp != 100:
                tekst = (
                    f"At each expiry date:\n"
                    f"• If spot is at or below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} {tek} Notional at {((strike - wyp / 100)):.{s_decimal}f}\n"
                    f"• If spot is above {trigger:.{s_decimal}f} and below {strike:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} {tek} Notional at {strike:.{s_decimal}f}\n"
                    f"• If spot is at or above {strike:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} Notional at {strike:.{s_decimal}f}"
                )
        if not is_importer and trigger >= strike:
            if zp == 100:
                tekst = (
                    f"At each expiry date:\n"
                    f"• If spot is at or above {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} Notional at {((strike - wyp / 100)):.{s_decimal}f}\n"
                    f"• If spot is below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} at {strike:.{s_decimal}f}"
                )
            elif zp != 100:
                tekst = (
                    f"At each expiry date:\n"
                    f"• If spot is at or above {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} {tek} Notional at {((strike - wyp / 100)):.{s_decimal}f}\n"
                    f"• If spot is above {strike:.{s_decimal}f} and below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {tek} {base} Notional at {strike:.{s_decimal}f}\n"
                    f"• If spot is at or below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} at {strike:.{s_decimal}f}"
                )
        #scenario 2 ___
        if is_importer and trigger >= strike:
            if zp == 100:
                tekst = (
                    f"At each expiry date: \n"
                    f"• If spot is at or below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} Notional at {strike - wyp / 100:.{s_decimal}f}\n"
                    f"• If spot is above {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} Notional at {strike:.{s_decimal}f}"
                )
            elif zp != 100:
                tekst = (
                    f"At each expiry date: \n"
                    f"• If spot is at or below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} {tek} Notional at {strike - wyp / 100:.{s_decimal}f}\n"
                    f"• If spot is above {strike:.{s_decimal}f} and below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} {tek} Notional at {strike - wyp / 100:.{s_decimal}f}\n"
                    f"• If spot is at or above {trigger:.{s_decimal}f} {company}, "
                    f"{company} will buy {base} Notional at {strike:.{s_decimal}f}"
                )
        if not is_importer and trigger <= strike:
            if zp == 100:
                tekst = (
                    f"At each expiry date: \n"
                    f"• If spot is at or above {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} Notional at {strike - wyp / 100:.{s_decimal}f}\n"
                    f"• If spot is below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} at {strike:.{s_decimal}f}"
                )
            elif zp != 100:
                tekst = (
                    f"At each expiry date: \n"
                    f"• If spot is at or above {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} {tek} Notional at {strike - wyp / 100:.{s_decimal}f}\n"
                    f"• If spot is above {strike:.{s_decimal}f} and below {trigger:.{s_decimal}f} {company}, "
                    f"{company} will sell {tek} {base} Notional at {strike:.{s_decimal}f}\n"
                    f"• If spot is at or below {strike:.{s_decimal}f} {company}, "
                    f"{company} will sell {base} at {strike:.{s_decimal}f}"
                )
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = f"• {base}/{term} Spot: {spot:.{s_decimal}f}\n" \
        f"• Expiry Date: {expiry_date}\n" \
        f"• Strike: {strike:.{s_decimal}f}\n" \
        f"• Trigger: {trigger:.{s_decimal}f} \n" \
        f"• Reset Strike: {strike - (wyp / 100):.{s_decimal}f} \n" \
        f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if lever!=1: terms+=f"\n• Leveraged Notional: {'{:,.0f}'.format(notional)}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Client knows that maximum {base} buying rate will not be higher than Strike\n• Possibility of improving {base} buying rate if spot at expiry trades below Trigger",
        f"• If spot at expiry fixes above Trigger there will be no improvement in the Strike\n• Opportunity cost if spot at expiry is below effective {base} buying rate"]
        if not is_importer: pros_cons = [f"• Client knows that minimum {base} selling rate will not be worse than Strike\n• Possibility of improving {base} selling rate if spot at expiry trades above Trigger",
        f"• If spot at expiry at expiry fixes below Trigger there will be no improvement in the Strike\n• Opportunity cost if spot at expiry is above effective {base} selling rate"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Forward with the added benefit of improving {base} buying rate if spot fixes below Trigger"
        if not is_importer: subtitle = f"Forward with the added benefit of improving {base} selling rate if spot fixes above Trigger"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    #scenario
        if lever==1: scenario = forw_zrek_scenario(strike, wyp, zp, trigger, notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = forw_zrek_scenario_lever(strike, wyp, zp,lever, trigger, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Leveraged Forwad" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Leveraged Forward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        zp = 1/float(row[1].get())*100
        lever = float(row[1].get())
        array_ = forward_fun(strike,zp,is_importer, c_min, c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 14
        else: slide_n = 42
    
    #----- fix! - pierwszy wiersz blad
    
    # payout
        if is_importer: tekst = f"At each expiry date:\n" \
    f"• If spot is at or below {strike:{s_decimal}f} {company} will buy {base} Leveraged Notional at {strike:{s_decimal}f}\n" \
    f"• If spot is above {strike:{s_decimal}f} {company} will buy {base} Notional at {strike:{s_decimal}f} and the rest at market spot rate\n" \
    f"• Reset Strike: {strike - (wyp / 100):.{s_decimal}f}\n" \
    f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if not is_importer: tekst =  f"At each expiry date:\n" \
    f"• If spot is above or at {strike:{s_decimal}f} {company} will sell {base} Leveraged Notional at {strike:{s_decimal}f}\n" \
    f"• If spot is below {strike:{s_decimal}f} {company} will sell {base} Notional at {strike:{s_decimal}f} and the rest at market spot rate\n" 
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms =f"• {base}/{term} Spot: {float(spot):.{s_decimal}f}\n" + \
       f"• Expiry Date: {expiry_date}\n" + \
       f"• Strike: {float(strike):.{s_decimal}f}\n" + \
       f"• Notional: {'{:,.0f}'.format(int(notional/lever))}"
        if lever!=1: terms+=f"\n• Leveraged Notional: {'{:,.0f}'.format(notional)}"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Ability to buy {base} below current Forward rate\n",
        f"• Client only has partial protection\n• No participation in {base} weakening\n• Obligation to buy {base} Leveraged Notional if spot at expiry is below Strike"]
        if not is_importer: pros_cons = [f"• Ability to sell {base} above current Forward rate",
        f"• Client only has partial protection\n• No participation in {base} strengthening\n• Obligation to sell {base} Leveraged Notional if spot at expiry is above Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Offers potential improvement in {base} buying rate, but not a full hedge"
        if not is_importer: subtitle = f"Offers potential improvement in {base} selling rate, but not a full hedge"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    # scenario
        scenario = forward_scenario(strike, zp, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Seagull" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Seagull"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        mid = float(row[1].get())
        upper = float(row[2].get())
        zp = 1/float(row[3].get())*100
        lever = float(row[3].get())
        array_ = kor_mewa_fun(lower, upper, mid, zp,
                           is_importer, c_min, c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 15
        else: slide_n = 43
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Allows to participate in {base} rate strengthening down to Lower Strike\n• Not a full hedge structure, but if spot goes above Upper Strike Client receives a subsidy",
        f"• No hedge above Upper Strike\n• Below Lower Strike Client will be obliged to buy {base} Notional\n• Opportunity cost if spot at expiry is below Lower Strike"]
        if not is_importer: pros_cons = [f"• Allows to participate in {base} rate strengthening up to Upper Strike\n• Not a full hedge structure, but if spot goes below Lower Strike Client receives a subsidy",
        f"• No hedge below Lower Strike\n• Above Upper Strike Client will be obliged to sell {base} Notional\n• Opportunity cost if spot at expiry is above Upper Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Allows to participate down to Lower Strike and above Upper Strike gives subsidy compared to market rate"
        if not is_importer: subtitle = f"Allows to participate up to Upper Strike and below Lower Strike gives subsidy compared to market rate"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    # scenario
        if lever==1: scenario = kor_mewa_scenario(lower, upper, mid, zp, notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = kor_mewa_scenario_lever(lower, upper, mid, zp,lever, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Lookback Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Lookback Forward"
        row = find_row(selected_inputs, s)
        reset = float(row[0].get())
        reference = float(row[1].get())
        range_min = float(row[2].get())
        range_max = float(row[3].get())
        adj = float(row[4].get())
        observation_date = row[5].get()
        expiry_date = row[6].get()
        if is_importer: slide_n = 16
        else: slide_n = 44
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Full protection against higher spot at {base} rate not worse than Reset Strike\n• Ability to lock in spot at the lowest fixing traded in the Range",
        f"• Reset Strike is higher compared to Forward rates\n• Opportunity cost if spot breaks the Range"]
        if not is_importer: pros_cons = [f"• Full protection against lower spot at {base} rate not worse than Reset Strike\n• Ability to lock in spot at the highest fixing traded in the Range",
        f"• Reset Strike is lower compared to Forward rates\n• Opportunity cost if spot breaks the Range"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Solution hedges {company} against {base} strengthening while allowing to lock in lower spot if {term} appreciates but stays in the pre-agreed Range over Observation Period"
        if not is_importer: subtitle = f"Solution hedges {company} against {base} weakening while allowing to lock in higher spot if {term} depreciates but stays in the pre-agreed Range over Observation Period"
        do_subtitle(presentation,subtitle,slide_n)
        do_eff(slide,base,is_importer)
        
    # scenario

####################################################################################################################################        
    
    if "Convertible Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Convertible Forward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        trigger = float(row[1].get())
        zp = 1/float(row[2].get())*100
        lever = float(row[2].get())
        array_ = forw_elas_fun(strike,trigger , zp,
                           is_importer, c_min, c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 18
        else: slide_n = 46
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Full hedge protection with participation in lower {base} buying rate up to Trigger",
        f"• Maximum {base} buying rate is higher compared to Forward\n• Opportunity cost if spot at expiry is at or below Trigger"]
        if not is_importer: pros_cons = [f"• Full hedge protection with participation in higher {base} selling rate up to Trigger",
        f"• Minimum {base} selling rate is lower compared to Forward\n• Opportunity cost if spot at expiry is at or above Trigger"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Full hedge structure which offers participation in favorable spot move"
        if not is_importer: subtitle = f"Full hedge structure which offers participation in favorable spot move"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        if lever==1: scenario = forw_elas_scenario(strike, zp, trigger, notional, is_importer, s_decimal, s_min, s_max)
        #fix
        if lever!=1: scenario = forw_elas_scenario(strike, zp, trigger, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Knock-Out Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Knock-Out Forward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        ko = float(row[1].get())
        zp = 1/float(row[2].get())*100
        lever = float(row[2].get())
        array_ = forw_zwyl_iwyp_fun(strike,ko,0,zp,is_importer,c_min,c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 19
        else: slide_n = 47
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• If spot at expiry is below Knock-Out, {base} buying rate is lower compared to Forward",
        f"• No hedge if spot at expiry fixes above Knock-Out\n• Opportunity cost if spot at expiry is below Strike"]
        if not is_importer: pros_cons = [f"• If spot at expiry is above Knock-Out, {base} selling rate is higher compared to Forward",
        f"• No hedge if spot at expiry fixes below Knock-Out\n• Opportunity cost if spot at expiry is at or above Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Structure with improved strike rate but no hedge if spot fixes above Knock-Out level"
        if not is_importer: subtitle = f"Structure with improved strike rate but no hedge if spot fixes below Knock-Out level"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
        
    # scenario
        if lever==1: scenario = forw_zwyl_scenario(strike, ko, zp, notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = forw_zwyl_scenario_lever(strike, ko, zp,lever, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Knock-Out Collar" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Knock-Out Collar"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        ko = float(row[2].get())
        zp = 1/float(row[3].get())*100
        lever = float(row[3].get())
        array_ = kor_zwyl_iwyp_fun(lower, upper,ko, 0, zp,is_importer,c_min,c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 20
        else: slide_n = 48
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Client has hedge up to Knock-out and participates in {base} weakening down to Lower Strike",
        f"• No hedge if spot at expiry fixes at or above Knock-Out\n• Opportunity cost if spot at expiry is below Lower Strike"]
        if not is_importer: pros_cons = [f"• Client has hedge down to Knock-out and participate in {base} strengthening up to Upper Strike",
        f"• No hedge if spot at expiry fixes at or below Knock-Out\n• Opportunity cost if spot at expiry is above Upper Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Structure with improved strike rate which offers participation in favorable spot move but no hedge if spot fixes above Knock-Out"
        if not is_importer: subtitle = f"Structure with improved strike rate which offers participation in favorable spot move but no hedge if spot fixes below Knock-Out"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
        
    # scenario
        if lever==1: scenario = kor_zwyl_iwyp_scenario(lower, upper, ko, 100, zp, notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = kor_zwyl_iwyp_scenario_lever(lower, upper, ko, 100, zp,lever, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Knock-Out Collar with Subsidy" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Knock-Out Collar with Subsidy"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        ko = float(row[2].get())
        wyp = float(row[3].get())
        zp = 1/float(row[4].get())*100
        lever = float(row[4].get())
        array_ = kor_zwyl_iwyp_fun(lower,upper,ko,wyp,zp,is_importer,c_min,c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 21
        else: slide_n = 49
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Client has hedge up to Trigger and can participate in {base} weakening down to Lower Strike\n• If spot is at or above Knock-Out, Client receives a subsidy",
        f"• No hedge if spot at expiry fixes at or above Knock-Out\n• Opportunity cost if spot at expiry is below Lower Strike"]
        if not is_importer: pros_cons = [f"• Client has hedge down to Trigger and can participate in {base} strengthening up to Upper Strike\n• If spot is at or below Knock-Out, Client receives a subsidy",
        f"• No hedge if spot at expiry fixes at or below Knock-Out\n• Opportunity cost if spot at expiry is above Upper Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Structure with improved strike rate which offers participation in favorable spot move. In case of losing hedge when spot fixes above Knock-Out Client receives subsidy"
        if not is_importer: subtitle = f"Structure with improved strike rate which offers participation in favorable spot move. In case of losing hedge when spot fixes below Knock-Out Client receives subsidy"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
        
    # scenario
        if lever==1: scenario = kor_zwyl_iwyp_scenario(lower, upper, ko, wyp, zp, notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = kor_zwyl_iwyp_scenario_lever(lower, upper, ko, wyp, zp,lever, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Knock-Out Forward with Subsidy" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Knock-Out Forward with Subsidy"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        ko = float(row[1].get())
        wyp = float(row[2].get())
        zp = 1/float(row[3].get())*100
        lever = float(row[3].get())
        array_ = forw_zwyl_iwyp_fun(strike,ko,0,zp,is_importer,c_min,c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 22
        else: slide_n = 50
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• If spot at expiry is below Knock-Out, {base} buying rate is lower compared to Forward\n• If spot is at or above Knock-Out, Client receives a subsidy",
        f"• No hedge if spot at expiry fixes at or above Knock-Out\n• Opportunity cost if spot at expiry is below Strike"]
        if not is_importer: pros_cons = [f"• If spot at expiry is above Knock-Out, {base} selling rate is higher compared to Forward\n• If spot is at or below Knock-Out, Client receives a subsidy",
        f"• No hedge if spot at expiry fixes at or below Knock-Out\n• Opportunity cost if spot at expiry is above Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitles
        if is_importer: subtitle = f"Structure with improved strike rate compared to Forward. In case of losing hedge when spot fixes above Knock-Out {company} receives subsidy"
        if not is_importer: subtitle = f"Structure with improved strike rate compared to Forward. In case of losing hedge when spot fixes below Knock-Out {company} receives subsidy"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        #fix
        if lever==1: scenario = forw_zwyl_scenario(strike, ko, 100, notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = forw_zwyl_iwyp_scenario_lever(strike, ko, 100,zp,lever, notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################
        
    if "Step Down Collar" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Step Down Collar"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        wyp = float(row[2].get())
        trigger = float(row[3].get())
        zp = 1/float(row[4].get())*100
        lever = float(row[4].get())
        if not is_importer: s = "Step Up Collar"
        array_ = kor_zrek_fun(lower,upper,trigger,wyp,zp,is_importer, c_min, c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 23
        else: slide_n = 51
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Client knows that maximum {base} buying rate will not be worse than Upper Strike\n• Possibility of improving {base} buying rate in case spot at expiry fixes below Trigger",
        f"• If spot at expiry fixes above Trigger, Lower Strike in the structure is not improved\n• Opportunity cost if spot at expiry is below effective {base} buying rate"]
        if not is_importer: pros_cons = [f"• Client knows that minimum {base} selling rate will not be worse than Lower Strike\n• Possibility of improving {base} selling rate in case spot at expiry fixes above Trigger",
        f"• If spot at expiry fixes below Trigger, Upper Strike in the structure is not improved\n• Opportunity cost if spot at expiry is above effective {base} selling rate"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Collar with the added benefit of improving {base} buying rate if spot fixes below Trigger"
        if not is_importer: subtitle = f"Collar with the added benefit of improving {base} selling rate if spot fixes above Trigger"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        if lever==1: scenario = Kor_zrek_scenario(lower,upper,zp,trigger,wyp,notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = kor_zrek_scenario_lever(lower,upper,zp,lever,trigger,wyp,notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Convertible Collar" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Convertible Collar"
        row = find_row(selected_inputs, s)
        lower = float(row[0].get())
        upper = float(row[1].get())
        trigger = float(row[2].get())
        zp = 1/float(row[3].get())*100
        lever = float(row[3].get())
        array_ = kor_elas_fun(lower,upper,trigger,zp,is_importer,c_min,c_max)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 24
        else: slide_n = 52
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Client knows the maximum {base} buying rate and can participate in lower spot down to Trigger\n• In case of {base} weakening below Trigger {base} buying rate is lower than Upper Strike",
        f"• If spot trades at or below Trigger, Client will be obliged to buy {base} at Lower Strike which will be higher than the prevailing spot market"]
        if not is_importer: pros_cons = [f"• Client knows the minum {base} selling rate and can participate in higher spot up to Trigger\n• In case of {base} strengthening above Trigger {base} selling rate is higher than Lower Strike",
        f"• If spot trades at or above Trigger, Client will be obliged to sell {base} at Upper Strike which will be lower than the prevailing spot market"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Full hedge structure which offers participation in favorable spot move"
        if not is_importer: subtitle = f"Full hedge structure which offers participation in favorable spot move"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        if lever==1: scenario = Kor_elas_scenario(lower,upper,zp,trigger,notional, is_importer, s_decimal, s_min, s_max)
        if lever!=1: scenario = Kor_elas_scenario_lever(lower,upper,zp,lever,trigger,notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Extendable Forward" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Extendable Forward"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        guranteed_date = row[1].get()
        conditional_date = row[2].get()
        zp = 1/float(row[3].get())*100
        lever = float(row[3].get())
        array_ = forward_fun(strike,zp,is_importer,c_min,c_max,c_decimal)
        chart_data = do_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 25
        else: slide_n = 53
    
   # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Lower {base} buying rate compared to Forward\n• On the Guaranteed Expiry Dates Client has full hedge in strips of Forward with known {base} buying rate",
        f"• If spot trades below Strike Client will be obliged to buy Notional at Strike which will be higher than market spot rate\n• If structure is extended Client will be obliged to buy rate at worse than market level if spot at expiry trades below Strike"]
        if not is_importer: pros_cons = [f"• Higher {base} selling rate compared to Forward\n• On the Guaranteed Expiry Dates Client has full hedge in strips of Forward with known {base} selling rate",
        f"• If spot goes above Strike Client will be obliged to sell Notional at Strike which will be lower than market spot rate\n• If structure is extended Client will be obliged to sell rate at worse than market level if spot at expiry trades above Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"{company} buys {base} at a better level compared to Forward by giving the right to extend the tenor of the hedge"
        if not is_importer: subtitle = f"{company} sells {base} at a better level compared to Forward by giving the right to extend the tenor of the hedge"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    # scenario
        scenario = forward_scenario(strike,zp,notional, is_importer, s_decimal, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "TARF" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
        s = "TARF"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        zp = float(row[4].get())
        lever = float(row[4].get())
        array_ = forward_fun(strike, zp, is_importer, c_min, c_max, c_decimal)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 26
        else: slide_n = 54
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Lower buying rate compared to Forward until Cap is reached\n• Trade will only terminate when\n• Client has accrued a pre-agreed profit\nOnce cap is reached Client can enter into new TARF",
        f"• Once the Cap is reached structure knocks out leaving Client unhedged\n• Levels in new TARF, if Client decides to enter one, might be worse than in initial TARF or Forward"]
        if not is_importer: pros_cons = [f"• Higher {base} selling rate compared to Forward\n• On the Guaranteed Expiry Dates Client has full hedge in strips of Forward with known {base} selling rate",
        f"• If spot goes above Strike Client will be obliged to sell Notional at Strike which will be lower than market spot rate\n• If structure is extended Client will be obliged to sell rate at worse than market level if spot at expiry trades above Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Solution for clients who would like to enhance their buying rate compared to vanilla Forwards. Structure performs best when spot trades in a narrow range. Client accepts the risk of being unhedged if structure gets knocked out"
        if not is_importer: subtitle = f"Solution for clients who would like to enhance their selling rate compared to vanilla Forwards. Structure performs best when spot trades in a narrow range. Client accepts the risk of being unhedged if structure gets knocked out"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
        do_eff(slide,base,is_importer)
    
    # scenarios
        subtitle = f"Scenario analysis presents different sample paths of spot fixing rate over the life of transaction showing example outcomes of TARF"
        do_subtitle(subtitle,slide_n+1)
        do_subtitle(subtitle,slide_n+2)
        
###########################################################################################################################################################################
    
    if "American Call Option" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "American Call Option"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        koszt = float(row[1].get())
        nominal = float(''.join(row[2].get().split()))
        if not is_importer: s = "American Put Option"
        array_ = Eur_fun(strike, 100, koszt, spot, is_importer, c_min, c_max)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 29
        else: slide_n = 57
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Possibility of exercising bought right (even partially) anytime from trade date till expiry date\n• Unlike Forward contract, option offers full participation in lower {base} buying rate\n• Full hedge above Strike level",
        f"• Upfront premium payment"]
        if not is_importer: pros_cons = [f"• Possibility of exercising sold right (even partially) anytime from trade date till expiry date\n• Unlike Forward contract, option offers full participation in higher {base} selling rate\n• Full hedge below Strike level",
        f"• Upfront premium payment"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in lower spot move with possibility of exercising bought right anytime from trade date till Expiry Date"
        if not is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in higher spot move with possibility of exercising bought right anytime from trade date till Expiry Date"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        scenario = Eur_scenario(strike, nominal, is_importer, s_decimal,zp, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "European Call Option with Rebate" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "European Call Option with Rebate"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        koszt = float(row[1].get())
        nominal = float(''.join(row[2].get().split()))
        reb_level = float(row[3].get())
        reb = float(row[4].get())
        if not is_importer: s = "European Put Option with Rebate"
        array_ = Eur_r_fun(strike,reb_level,reb, 100, koszt, spot, is_importer, c_min, c_max)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 30
        else: slide_n = 58
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Premium is partially paid back to Client if rate at given expiry date trades below Rebate rate\n• Unlike Forward contract, option offers full participation in lower {base} buying rate\n• Full hedge above Strike level",
        f"• Upfront premium payment\n• Higher cost compared to buying Call Option without rebate"]
        if not is_importer: pros_cons = [f"• Premium is partially paid back to Client if rate at given expiry date trades above Rebate rate\n• Unlike Forward contract, option offers full participation in higher {base} selling rate\n• Full hedge below Strike level",
         f"• Upfront premium payment\n• Higher cost compared to buying Put Option without rebate"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in lower spot move and possibility of receiving back paid premium provided spot moves below predefined rate"
        if not is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in higher rate spot move and possibility of receiving back paid premium provided spot moves below predefined rate"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        scenario = Eur_r_scenario(strike,reb_level, reb,nominal, is_importer, s_decimal,100, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "American Call Option with Rebate" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "American Call Option with Rebate"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        koszt = float(row[1].get())
        nominal = float(''.join(row[2].get().split()))
        reb_level = float(row[2].get())
        reb = float(row[3].get())
        if not is_importer: s = "American Put Option with Rebate"
        array_ = Eur_r_fun(strike,reb_level,reb, 100, koszt, spot, is_importer, c_min, c_max)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 31
        else: slide_n = 59
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Premium is partially paid back to Client if {base} rate at given expiry date trades below Rebate rate\n• Possibility of exercising bought right anytime from trade date till expiry date, even partially\n• Full hedge above Strike level",
        f"• Upfront premium payment\n• Higher cost compared to buying American Call Option without rebate"]
        if not is_importer: pros_cons = [f"• Premium is partially paid back to Client if {base} rate at given expiry date trades above Rebate rate\n• Possibility of exercising sold right anytime from trade date till expiry date, even partially\n• Full hedge below Strike level",
        f"• Upfront premium payment\n• Higher cost compared to buying American Put Option without rebate"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in lower spot move. Company has the right to exercise option anytime from trade date till Expiry Date and also possibility of getting back part of the paid premium."
        if not is_importer: subtitle = f"Provides a hedge at the strike level while retaining full participation in higher rate spot move and possibility of receiving back paid premium and possibility of being exercised anytime from trade date till Expiry Date"
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        scenario = Eur_r_scenario(strike,reb_level, reb,nominal, is_importer, s_decimal,100, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "Option on Option" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
    # data
        s = "Option on Option"
        row = find_row(selected_inputs, s)
        strike = float(row[0].get())
        koszt = float(row[1].get())
        nominal = float(''.join(row[2].get().split()))
        array_ = Eur_fun(strike, 100, koszt, spot, is_importer, c_min, c_max)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        if is_importer: slide_n = 32
        else: slide_n = 60
    
    # payout
        if is_importer: tekst = ""
        if not is_importer: tekst = ""
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = ""
        do_terms(presentation,terms, slide_n)
    # pros/cons
        if is_importer: pros_cons = [f"• Full hedge above Strike level\n• Unlike forward contract, option offers full participation in {base} lower rate\n• Client can split premium into 2 payments and can walk away from the transaction if option is not needed on the Decision Date",
        f"• Total premium cost, if both Initial and Future Premiums are paid, is higher compared to Vanilia Option with same Expiry Date and Strike"]
        if not is_importer: pros_cons = [f"• Full hedge below Strike level\n• Unlike forward contract, option offers full participation in {base} stronger rate\n• Client can split premium into 2 payments and can walk away from the transaction if option is not needed on the Decision Date",
        f"• Total premium cost, if both Initial and Future Premiums are paid, is higher compared to Vanilia Option with same Expiry Date and Strike"]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        if is_importer: subtitle = f"Provides full hedge at the strike level while retaining participation in lower spot moves. At trade time Client pays Initial Premium and if on Decision Date option is still needed, additional premium is paid."
        if not is_importer: subtitle = f"Provides full hedge at the strike level while retaining participation in lower spot moves. At trade time Client pays Initial Premium and if on Decision Date option is still needed, additional premium is paid."
        do_subtitle(presentation,subtitle,slide_n)
        
    # chart
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        chart_data.close()
    
    # scenario
        scenario = Eur_scenario(strike, nominal, is_importer, s_decimal,zp, s_min, s_max)
        put_table_on_slide(scenario,s_decimal,slide, header_name)
        do_eff(slide,base,is_importer)

####################################################################################################################################

    if "New Instrument" in selected_options:
        counter+=1
        loading_label.config(text=f"Loading: {int(counter/len(selected_options)*100)}%")
        loading_label.update()
        s = "New Instrument"
        slide_n = 61
    
    # payout
        tekst = "Payout for Your Instrument"
        do_tekst(presentation,tekst, slide_n)
    # why to use    
        if is_importer: whyto = ""
        if not is_importer: whyto = ""
    # terms
        terms = "Terms for Your Instrument"
        do_terms(presentation,terms, slide_n)
    # pros/cons
        pros_cons = ["Pros...","Cons..."]
        do_pros_cons(presentation,pros_cons,slide_n)
    # subtitle
        subtitle = f"Subtitle for Your Instrument"
        do_subtitle(presentation,subtitle,slide_n)
        do_eff(slide,base,is_importer)
        
    # chart

        """
        array_ = Eur_fun(strike, 100, koszt, spot, is_importer, c_min, c_max)
        chart_data = do_option_chart(s,base,term,c_min,c_max,c_decimal,array_)
        
        slide_n = 55
        slide = presentation.slides[slide_n]
        image_stream = plot_to_image(chart_data)
        slide.shapes.add_picture(image_stream, left=Inches(to_cm(1.95)), top=Inches(to_cm(8.49)), 
                                 width=Inches(to_cm(10)), height=Inches(to_cm(6.6)))
        """

    loading_label.config(text=f"Ready")
    presentation.save("rms_presentation.pptx")

####################################################################################################################################
#
# help functions
#
####################################################################################################################################

def do_tekst(presentation,tekst, slide_n):
    slide = presentation.slides[slide_n]
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(13.88)), top=Inches(to_cm(4.36)), width=Inches(to_cm(13.02)), height=Inches(to_cm(3.78)))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = tekst
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Citi Sans Text"
        paragraph.space_after = Pt(3) 

def do_terms(presentation, terms, slide_n):
    slide = presentation.slides[slide_n]
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(0.63)), top=Inches(to_cm(4.36)), width=Inches(to_cm(7.53)), height=Inches(to_cm(3.71)))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = terms
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Citi Sans Text"
        paragraph.space_after = Pt(3) 
    
def do_pros_cons(presentation,pros_cons, slide_n):
    slide = presentation.slides[slide_n] # PROS
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(3.08)), top=Inches(to_cm(15.8)), width=Inches(to_cm(9.25)), height=Inches(to_cm(1.64)))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = pros_cons[0]
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Citi Sans Text"
        paragraph.space_after = Pt(3) 
    
    slide = presentation.slides[slide_n] # CONS
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(16.53)), top=Inches(to_cm(15.8)), width=Inches(to_cm(9.25)), height=Inches(to_cm(1.64)))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = pros_cons[1]
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Citi Sans Text"
        paragraph.space_after = Pt(3)
   
def do_subtitle(presentation,subtitle, slide_n):
    slide = presentation.slides[slide_n]
    textbox = slide.shapes.add_textbox(left=Inches(to_cm(1.1)), top=Inches(to_cm(2.15)), width=Inches(to_cm(25.22)), height=Inches(to_cm(1.8)))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = subtitle
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Citi Sans Text"
        paragraph.space_after = Pt(3) 
    

def to_cm(x): return x/2.54
def find_row(array, key):
    for row in array:
        if row[0] == key:
            return row[1:]
        
def plot_to_image(plt):
    image_stream = BytesIO()
    plt.savefig(image_stream, format='png', bbox_inches='tight')
    image_stream.seek(0)
    return image_stream

def put_table_on_slide(scenario,s_decimal,slide, headers_name,partyp=False):
        if partyp==False: scenario = [[row[0], row[1], row[4]] for row in scenario]
        else: scenario = [[row[0], row[4], row[3]] for row in scenario]
        table = slide.shapes.add_table(rows=len(scenario)+1,cols=len(scenario[0]), left=Inches(to_cm(15.11)), 
                                       top=Inches(to_cm(8.72)), width=Inches(to_cm(8.73)), height=Inches(to_cm(5.45))).table
        for i, row in enumerate(scenario):
            for j, value in enumerate(row):
                cell = table.cell(i+1, j)
                cell.text = str(value)

        for j, header_text in enumerate(headers_name):
            cell = table.cell(0, j)  
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 92, 226)  
            cell.text = header_text
            cell.text_frame.paragraphs[0].font.size = Pt(9)  
            cell.text_frame.paragraphs[0].font.name = 'Arial'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
            cell.text_frame.paragraphs[0].font.bold = True  

        for i in range(0, len(table.rows)-1):
            for j in range(len(table.columns)):
                value = float(table.cell(i + 1, j).text)
                if j in (0, 1): format_cell_decimal(table.cell(i + 1, j), value, decimal_places=s_decimal)
                elif j == 2: format_cell_thousands(table.cell(i + 1, j), value)
        apply_formatting_to_table(table)

def format_cell_decimal(cell, value, decimal_places):
    cell.text = str(value)
    cell.text = f'{value:.{decimal_places}f}'
def format_cell_thousands(cell, value):
    cell.text = str(value)
    cell.text = '{:,.0f}'.format(value)

def do_eff(slide, base, is_importer):
    if is_importer:
        tekst1 = f"Effective {base} buying rate"
    else:
        tekst1 = f"Effective {base} selling rate"

    textbox = slide.shapes.add_textbox(Inches(to_cm(0.63)), Inches(to_cm(7.62)),
                                    Inches(to_cm(10.51)), Inches(to_cm(1.6)))
    text_frame = textbox.text_frame
    text_frame.text = tekst1
    text_frame.paragraphs[0].font.size = Pt(16)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(36, 91, 226)
    text_frame.paragraphs[0].font.name = "Citi Sans Display"
    text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

def apply_formatting_to_table(table):
    for i in range(1, len(table.rows)):
        for j in range(len(table.columns)):
            cell = table.cell(i, j)
            value = float(cell.text.replace(',', ''))
            if value < 0:
                for k in range(len(table.columns)):
                    cell = table.cell(i, k)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0) 
                    cell.text_frame.paragraphs[0].font.size = Pt(10) 
            else:
                for k in range(len(table.columns)):
                    cell = table.cell(i, k)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 92, 226) 
                    cell.text_frame.paragraphs[0].font.size = Pt(10) 
